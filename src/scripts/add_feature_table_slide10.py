#!/usr/bin/env python3
"""
Add binary classification SFS feature selection table to slide 10 of
UKB-DRP_完整实验报告.pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PPTX = "docs/PPT/UKB-DRP_完整实验报告.pptx"
BLUE = RGBColor(0x00, 0x66, 0xCC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Helvetica"


def set_font(run, size_pt, bold=False, color=GRAY):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def add_table(slide, left_in, top_in, col_widths_in, headers, rows):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths_in)
    row_h = 0.22

    shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left_in), Inches(top_in),
        Inches(total_w), Inches(row_h * n_rows),
    )
    table = shape.table
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    # Header
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(hdr)
        set_font(run, 7, bold=True, color=WHITE)
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': '0066CC'})
        sf.append(sc)
        tcPr.append(sf)

    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            set_font(run, 7, bold=(j == 0), color=DARK)
            if i % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F5F7FA'})
                sf.append(sc)
                tcPr.append(sf)

    return shape


def main():
    prs = Presentation(PPTX)
    slide = prs.slides[9]  # slide 10

    # ── Step 1: Move existing bottom tables down ──
    # Move the AUC results table (Table 6) from top=3.5 to top=4.35
    # Move the imaging increment table (Table 7) from top=3.5 to top=4.35
    # Move conclusion (TextBox 8) from top=5.85 to top=6.35
    # Move page number (TextBox 9) - keep as is

    for shape in slide.shapes:
        if shape.name == "Table 6":
            shape.top = Inches(4.35)
        elif shape.name == "Table 7":
            shape.top = Inches(4.35)
        elif shape.name == "TextBox 8":
            shape.top = Inches(6.35)

    # ── Step 2: Insert feature selection table ──
    # Position: left=0.6, top=2.65, width ~12.0, height ~1.5
    # Compact format: 4 rows (3yr/5yr/10yr/all-time) showing top features

    headers = ["窗口",
               "#1 特征 (+AUC)", "#2 特征 (+AUC)", "#3 特征 (+AUC)",
               "#4 特征 (+AUC)", "#5 特征 (+AUC)"]

    rows = [
        ["3yr",
         "FS_ST40TA (+0.723)", "AMY_SUMRY_SUVR (+0.770)",
         "FS_ST44TA (+0.788)", "FS_ST29SV (+0.796)",
         "FS_ST55TS (+0.808)"],
        ["5yr",
         "AMY_CENTILOIDS (+0.750)", "FS_ST24TA (+0.802)",
         "FS_ST84CV (+0.815)", "AMY_COMP_SUVR (+0.840)",
         "FS_ST103TA (+0.846)"],
        ["10yr",
         "NfL_F (+0.741)", "FS_ST31TA (+0.845)",
         "AMY_CENTILOIDS (+0.865)", "AB42_F (+0.900)",
         "WMH_CER_TCB (+0.907)"],
        ["all-time",
         "PTWORKHS (+0.728)", "FS_ST12SV (+0.775)",
         "AMY_CENTILOIDS (+0.775)", "FS_ST129SA (+0.817)",
         "FS_ST72SA (+0.826)"],
    ]

    col_widths = [0.75, 2.15, 2.15, 2.15, 2.15, 2.15]
    add_table(slide, 0.6, 2.65, col_widths, headers, rows)

    # ── Step 3: Add a small annotation label above the feature table ──
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(2.45), Inches(6.0), Inches(0.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "▎SFS 特征选择 — Top 5 每窗口 (Model A: Bio+Img)"
    set_font(run, 9, bold=True, color=BLUE)

    # ── Step 4: Update page number if needed ──
    # (keep existing)

    prs.save(PPTX)
    print(f"✓ Updated {PPTX}")
    print(f"  Added SFS feature selection table to slide 10")


if __name__ == "__main__":
    main()
