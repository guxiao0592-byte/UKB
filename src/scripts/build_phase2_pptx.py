#!/usr/bin/env python3
"""
Build Phase 2 slides and append to existing UKB-DRP presentation.

Usage:
    source enve/bin/activate
    python src/scripts/build_phase2_pptx.py
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import csv

# ── Paths ──────────────────────────────────────────────────
PPTX_SRC = "docs/PPT/UKB-DRP_完整实验报告_副本.pptx"
PPTX_OUT = "docs/PPT/UKB-DRP_完整实验报告_副本.pptx"  # overwrite
RESULTS_DIR = "local_data/Results_adni"

# ── Style Constants (from existing PPTX analysis) ──────────
BLUE = RGBColor(0x00, 0x66, 0xCC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_FILL = RGBColor(0xE8, 0xF0, 0xFE)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
FONT_FAMILY = "Helvetica"

# ── Helper Functions ───────────────────────────────────────

def _set_font(run, size_pt, bold=False, color=GRAY, name=FONT_FAMILY):
    """Set font properties on a run."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_blue_accent_bar(slide, top_in=0.0, height_in=0.06):
    """Full-width thin blue bar at top of slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top_in), SLIDE_W, Inches(height_in),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return bar


def add_title_block(slide, title_text, subtitle_text, top_in=0.25):
    """Standard title + subtitle block."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_in), Inches(11.5), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_font(run, 25, bold=True, color=DARK)

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(top_in + 0.45), Inches(11.5), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = subtitle_text
    _set_font(run2, 12, bold=False, color=GRAY)


def add_page_number(slide, page_num):
    """Bottom-right page number."""
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num}/33"
    _set_font(run, 7, bold=False, color=GRAY)


def add_card(slide, left_in, top_in, width_in, height_in, number, heading, body_lines):
    """Standard info card with number, heading, and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.05)

    # Number
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(number)
    _set_font(run, 24, bold=True, color=BLUE)

    # Heading
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)
    run2 = p2.add_run()
    run2.text = heading
    _set_font(run2, 11, bold=True, color=DARK)

    # Body
    for line in body_lines:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        p3.space_before = Pt(1)
        run3 = p3.add_run()
        run3.text = line
        _set_font(run3, 9, bold=False, color=GRAY)

    return shape


def add_table(slide, left_in, top_in, col_widths_in, headers, rows,
              header_bg=BLUE, header_fg=WHITE):
    """Create a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths_in)

    shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left_in), Inches(top_in),
        Inches(total_w), Inches(0.35 * n_rows),
    )
    table = shape.table

    # Set column widths
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    # Header row
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(hdr)
        _set_font(run, 9, bold=True, color=header_fg)
        # Header background
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
        srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': '0066CC'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            _set_font(run, 9, bold=False, color=DARK)
            # Alternating row background
            if i % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
                srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F5F7FA'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

    return shape


def add_section_divider(slide, part_label, title, subtitle, page_num):
    """Section divider slide matching Slides 3/9 style."""
    # Full-width blue bars top and bottom
    add_blue_accent_bar(slide, 0.0, 0.08)
    add_blue_accent_bar(slide, 7.28, 0.08)

    # Part label
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = part_label
    _set_font(run, 18, bold=True, color=BLUE)

    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = title
    _set_font(run2, 32, bold=True, color=DARK)

    # Subtitle
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = subtitle
    _set_font(run3, 14, bold=False, color=GRAY)

    add_page_number(slide, page_num)


def add_bottom_conclusion(slide, top_in, text_lines, width_in=11.5):
    """Bottom conclusion bar."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_in), Inches(width_in), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (label, body) in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label
        _set_font(run, 11, bold=True, color=BLUE)
        if body:
            run2 = p.add_run()
            run2.text = f"  {body}"
            _set_font(run2, 10, bold=False, color=DARK)
    return txBox


# ── Slide Builders ─────────────────────────────────────────

def build_slide27_section(prs):
    """Slide 27: PART 4 section divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_section_divider(
        slide,
        part_label="PART 4",
        title="ADNI Phase 2 — 多终点预测",
        subtitle="DXSUM 二分类 → 生存分析 → CDR 量表扩展  |  MCI→Dementia + CN→MCI 双队列  |  s01–s05 统一管线",
        page_num=27,
    )


def build_slide28_overview(prs):
    """Slide 28: Phase 2 experiment overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2 实验总览 — 三阶段演进脉络",
                    "Phase 2A (DXSUM 二分类) → Phase 2B (DXSUM 生存) → Phase 2C (CDR 扩展 二分类+生存)")

    # 3-phase pipeline cards
    phases = [
        ("2A", "DXSUM 二分类",
         ["终点: MCI→Dementia (Broad)",
          "标签: 3yr/5yr/10yr 固定窗口",
          "模型: LightGBM + Isotonic",
          "删失浪费: 46–66% MCI 丢弃",
          "358 特征 → SFS Top 10"]),
        ("2B", "DXSUM 生存分析",
         ["终点: MCI→Dementia (time,event)",
          "标签: 连续时间, 保留全部 944 人",
          "模型: RSF + Cox PH",
          "输出: 完整 S(t|X) 生存曲线",
          "C-index 0.765, 一模型替代三窗口"]),
        ("2C", "CDR 量表扩展",
         ["终点: CDR 恶化 (4 种定义)",
          "队列: CN (633) + MCI (962)",
          "双任务: 二分类 + 生存",
          "CN→MCI 首次可预测 (C=0.68)",
          "CDR MCI 性能略超 DXSUM"]),
    ]
    for i, (label, title, lines) in enumerate(phases):
        add_card(slide, 0.6 + i * 4.15, 1.45, 3.9, 2.4, label, title, lines)

    # Arrow connectors
    for i in range(2):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(4.55 + i * 4.15), Inches(2.35), Inches(0.35), Inches(0.25),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLUE
        arrow.line.fill.background()

    # Bottom: unified methodology
    add_bottom_conclusion(slide, 4.15, [
        ("统一方法论 (s01–s05):", ""),
        ("s01", " LightGBM 5-fold Gain 排序 → Top 50  |  "),
        ("s02", " Ward 聚类 (|ρ|=0.75) 去冗余  |  "),
        ("s03", " 聚类后重排  |  "),
        ("s04", " SFS (Δ<0.0005 早停)  |  "),
        ("s05", " 5-fold CV 最终评估 (LGBM / RSF / Cox PH)"),
    ])

    add_page_number(slide, 28)


def build_slide29_phase2a(prs):
    """Slide 29: Phase 2A DXSUM binary classification labels & results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2A — DXSUM 二分类: 标签设计与结果",
                    "MCI→Dementia 三窗口 (3yr/5yr/10yr) + all-time  |  Model A (Bio+Img) vs Model B (Bio only)")

    # Left: label design table
    add_table(slide, 0.6, 1.4,
              [1.8, 1.3, 1.3, 1.3, 1.5, 1.2],
              ["窗口", "训练集 N", "事件 (y=1)", "阴性 (y=0)", "删失(丢弃)", "事件率"],
              [["3yr",  "600", "196", "404", "520 (46%)", "32.7%"],
               ["5yr",  "500", "267", "233", "620 (55%)", "53.4%"],
               ["10yr", "385", "311",  "74", "735 (66%)", "80.8%"],
               ["all-time", "1,120", "323", "797", "0", "28.8%"]])

    # Right: label rule card
    add_card(slide, 7.6, 1.4, 5.2, 1.8, "📋", "删失规则",
             ["y=1: converted AND time ≤ W 年",
              "y=0: NOT converted AND follow-up ≥ W 年",
              "删失: NOT converted AND follow-up < W 年 → 丢弃",
              "核心: 仅\"未转化+随访不够长\"才被丢弃"])

    # Bottom left: Model A results
    add_table(slide, 0.6, 3.5,
              [1.5, 1.0, 1.0, 1.15, 1.15, 1.0],
              ["窗口", "AUC", "Brier", "Sens", "Spec", "AUPRC"],
              [["3yr", "0.834±0.032", "0.154", "0.719", "0.795", "0.713"],
               ["5yr", "0.855±0.021", "0.161", "0.757", "0.798", "0.873"],
               ["10yr", "0.919±0.015", "0.087", "0.752", "0.959", "0.977"],
               ["all-time", "0.841±0.023", "0.144", "0.839", "0.686", "0.662"]])

    # Bottom right: Imaging increment + feature finding
    add_table(slide, 7.6, 3.5,
              [1.4, 1.3, 1.3, 1.2],
              ["窗口", "Bio only", "Bio+Img", "ΔAUC"],
              [["3yr", "0.730", "0.834", "+0.103 ← 最大"],
               ["5yr", "0.792", "0.855", "+0.063"],
               ["10yr", "0.825", "0.919", "+0.094"],
               ["all-time", "0.797", "0.841", "+0.044 ← 最小"]])

    add_bottom_conclusion(slide, 5.85, [
        ("核心发现:", "3yr 前6步全是影像 → 影像主导短期预测; NfL 是10yr #1 → 血浆标志物主导长期; PTWORKHS 是all-time #1"),
    ])

    add_page_number(slide, 29)


def build_slide30_survival_model(prs):
    """Slide 30: Phase 2B survival label design & model architecture changes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2B — 生存标签与模型架构调整",
                    "从固定窗口二分类 → 连续时间生存建模  |  (surv_time, surv_event) 一对标签替代三窗口")

    # Binary vs Survival comparison cards
    add_card(slide, 0.6, 1.4, 5.8, 2.6, "①", "二分类局限 (Phase 2A)",
             ["1. 删失浪费: 3yr丢46%, 5yr丢55%, 10yr丢66%",
              "2. 时间坍缩: 1.5年转化=2.9年转化=都是y=1",
              "3. 不一致风险: P(≤1yr)>P(≤3yr) 可能发生",
              "4. 回答不了\"什么时候转\": 只能固定窗口",
              "5. 需3个独立模型, 各自有不同训练集"])

    add_card(slide, 6.8, 1.4, 5.8, 2.6, "②", "生存标签优势 (Phase 2B)",
             ["1. 保留全部样本: 944人 (排除176零随访)",
              "  事件323人(34.2%) + 删失621人(65.8%)",
              "2. 连续时间: surv_time=转化年 or 最后随访年",
              "3. 自动一致: S(1)≥S(3)≥S(5) 数学保证",
              "4. 一个模型→任意时间点风险: risk_t = 1−S(t)",
              "5. 删失者贡献部分信息: P(T>t | censored)"])

    # Model architecture changes table
    add_table(slide, 0.6, 4.3,
              [2.0, 2.8, 2.8, 2.8, 2.0],
              ["组件", "二分类 (Phase 2A)", "v1 (C-index SFS)", "v2 (LGB AUC SFS)", "v2 优势"],
              [["标签 y", "Dementia_Wyr ∈ {0,1}", "(time, event)", "(time, event)", "—"],
               ["训练集", "按窗口筛选: 600/500/385", "统一 944 人", "统一 944 人", "—"],
               ["s01 排序", "LGBM Gain (分类)", "单变量 C-index", "LGBM Gain ★统一", "与2A一致"],
               ["s04 SFS", "AUC gain (分类)", "RSF C-index", "LGBM AUC ★统一", "与2A一致"],
               ["s05 主模型", "LGBMClassifier", "RSF (n=200,d=5)", "RSF (n=200,d=5)", "—"],
               ["s05 基准", "—", "Cox PH", "Cox PH (pen=0.1)", "—"],
               ["评估", "AUC, Brier, Sens", "C-index, tAUC", "C-index, tAUC, IBS", "—"]])

    add_page_number(slide, 30)


def build_slide31_survival_features(prs):
    """Slide 31: Phase 2B survival feature selection & ablation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2B — 生存特征选择与消融分析",
                    "v2 SFS: 3 Bio + 7 MRI  |  PTWORKHS 是不可替代的最强特征  |  Cox PH 风险比")

    # Left: SFS Top 10
    add_table(slide, 0.6, 1.4,
              [0.4, 2.2, 1.0, 1.0, 1.6],
              ["#", "特征", "AUC", "增益", "类型"],
              [["1", "PTWORKHS", "0.708", "+0.708", "Bio (就业状态)"],
               ["2", "FS_ST89SV", "0.755", "+0.047", "MRI 皮层下体积"],
               ["3", "AMY_CENTILOIDS", "0.756", "+0.001", "Amyloid PET"],
               ["4", "FS_ST72CV", "0.786", "+0.030", "MRI 皮层下体积"],
               ["5", "GFAP_Q", "0.795", "+0.009", "Bio (血浆 GFAP)"],
               ["6", "APOE4_count", "0.805", "+0.010", "Bio (遗传)"],
               ["7", "FS_ST93TA", "0.805", "−0.001", "MRI 皮层厚度"],
               ["8", "FS_ST12SV", "0.807", "+0.003", "MRI 皮层下体积"],
               ["9", "FS_ST102CV", "0.813", "+0.006", "MRI 皮层下体积"],
               ["10", "FS_ST84TS", "0.817", "+0.004", "MRI 厚度 SD"]])

    # Middle: Ablation waterfall
    add_table(slide, 4.5, 1.4,
              [1.0, 1.5, 1.1, 1.4],
              ["N特征", "C-index", "移除", "Δ"],
              [["10", "0.766", "(全)", "—"],
               ["9", "0.758", "FS_ST102CV", "−0.008"],
               ["6", "0.755", "APOE4", "−0.003"],
               ["5", "0.740", "GFAP_Q", "−0.015"],
               ["4", "0.740", "FS_ST72CV", "0"],
               ["3", "0.722", "AMY_CENTILOIDS", "−0.018"],
               ["2", "0.698", "FS_ST89SV", "−0.023"],
               ["1", "0.657", "PTWORKHS", "−0.082 ⚠"]])

    # Right: Cox PH hazard ratios
    add_table(slide, 8.5, 1.4,
              [2.0, 0.9, 1.0, 0.8],
              ["特征", "HR", "p值", "方向"],
              [["FS_ST12SV", "0.711", "1.5e-09 ***", "保护 ↓"],
               ["PTWORKHS", "0.737", "8.4e-08 ***", "保护 ↓"],
               ["FS_ST72CV", "0.772", "1.7e-06 ***", "保护 ↓"],
               ["FS_ST89SV", "1.238", "3.9e-06 ***", "风险 ↑"],
               ["APOE4_count", "1.242", "8.5e-06 ***", "风险 ↑"],
               ["FS_ST93TA", "0.842", "0.0005 ***", "保护 ↓"],
               ["AMY_CENTILOIDS", "1.138", "0.017 *", "风险 ↑"],
               ["FS_ST84TS", "0.890", "0.020 *", "保护 ↓"],
               ["GFAP_Q", "1.120", "0.026 *", "风险 ↑"],
               ["FS_ST102CV", "1.073", "0.180 NS", "—"]])

    add_bottom_conclusion(slide, 5.2, [
        ("核心发现:", "PTWORKHS (就业状态) 是最不可替代的特征 — 移除导致 C-index 暴跌 0.082。3 Bio (就业+GFAP+APOE4) + 7 MRI 构成最优组合。"),
    ])

    # Bottom comparison table
    add_table(slide, 0.6, 5.7,
              [2.2, 1.8, 1.8, 1.8, 1.8, 1.8],
              ["模型/指标", "C-index", "tAUC@1yr", "tAUC@3yr", "tAUC@5yr", "Brier@5yr"],
              [["RSF v2", "0.765±0.009", "0.753", "0.803", "0.834", "0.231"],
               ["Cox PH v2", "0.770", "—", "—", "—", "—"],
               ["RSF v1 (对比)", "0.745", "0.730", "0.791", "0.813", "0.277"],
               ["二分类 LGBM", "—", "—", "0.834", "0.855", "—"]])

    add_page_number(slide, 31)


def build_slide32_cdr_labels_features(prs):
    """Slide 32: Phase 2C CDR extension labels & feature comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2C — CDR 量表扩展: 终点定义与特征选择",
                    "4 种终点定义 × 2 队列 (CN+MCI)  |  CDR Primary/Sustained/CDRSB≥1/Composite")

    # Top: CDR label definitions
    add_card(slide, 0.6, 1.35, 5.8, 1.7, "📋", "CDR 终点定义",
             ["Primary: CDGLOBAL 首次恶化 (分层阈值)",
              "  CN (CDR=0): CDGLOBAL≥0.5  |  MCI (CDR=0.5): CDGLOBAL≥1.0",
              "Sustained: 首次恶化+下次访视确认 (降噪)",
              "CDRSB≥1: CDR Sum of Boxes 增加≥1.0分 (量化)",
              "Composite: CDGLOBAL恶化 OR DXSUM Dementia (仅MCI)"])

    # Right: CONSORT flow
    add_card(slide, 6.8, 1.35, 5.8, 1.7, "👥", "CONSORT 样本流 (1,595 纳入)",
             ["基线特征表: 2,749 → 有CDR数据: 1,595",
              "  CN 队列 (CDR=0→≥0.5): 633人, events=208 (32.9%)",
              "  MCI队列 (CDR=0.5→≥1): 962人, events=373 (38.8%)",
              "排除者更年轻 (−2.4yr) 且女性更多 (+11.6%)",
              "  → 存在选择偏差, 需在论文中说明"])

    # MCI Primary 3yr SFS
    add_table(slide, 0.6, 3.3,
              [0.4, 2.2, 1.0, 1.0, 1.6],
              ["#", "MCI 3yr SFS (AUC)", "AUC", "增益", "类型"],
              [["1", "FS_ST90TA", "0.709", "+0.709", "MRI 皮层厚度"],
               ["2", "FS_ST29SV", "0.782", "+0.073", "MRI 皮层下体积"],
               ["3", "FS_ST30SV", "0.802", "+0.020", "MRI 皮层下体积"],
               ["4", "pT217_AB42_F", "0.832", "+0.030", "血浆 pTau217/Aβ42"],
               ["5", "APOE4_count", "0.837", "+0.005", "遗传"],
               ["6", "FS_ST118CV", "0.847", "+0.010", "MRI 皮层下体积"],
               ["7", "FS_ST34TS", "0.858", "+0.011", "MRI 厚度SD"],
               ["8", "FS_ST40CV", "0.858", "+0.001", "MRI 皮层下体积"],
               ["9", "AMY_SUMMARY_SUVR", "0.862", "+0.004", "Amyloid PET"],
               ["10", "FS_ST85SA", "0.859", "−0.003", "MRI 表面积"]])

    # CN Primary 3yr SFS
    add_table(slide, 5.2, 3.3,
              [0.4, 2.4, 1.0, 1.0, 1.6],
              ["#", "CN 3yr SFS (AUC)", "AUC", "增益", "类型"],
              [["1", "TAU_META_TEMP_SUVR", "0.612", "+0.612", "Tau PET ★"],
               ["2", "FS_ST55CV", "0.649", "+0.037", "MRI 皮层下体积"],
               ["3", "AMY_COMP_REF_SUVR", "0.673", "+0.024", "Amyloid PET"],
               ["4", "FS_ST7SV", "0.671", "−0.003", "MRI 皮层下体积"],
               ["5", "FS_ST62SA", "0.687", "+0.017", "MRI 表面积"],
               ["6", "FS_ST119SA", "0.695", "+0.008", "MRI 表面积"],
               ["7", "FS_ST60TA", "0.705", "+0.010", "MRI 皮层厚度"],
               ["8", "FS_ST84TA", "0.721", "+0.016", "MRI 皮层厚度"],
               ["9", "FS_ST46TS", "0.723", "+0.002", "MRI 厚度SD"],
               ["10", "FS_ST102SA", "0.735", "+0.012", "MRI 表面积"]])

    add_bottom_conclusion(slide, 6.9, [
        ("跨队列差异:", "CN #1 = Tau PET (内嗅皮层), MCI #1 = MRI — Tau 在极早期已有独立预测信号。CN 前10步无Bio标志物入选, MCI 第4步即引入pTau217。"),
    ])

    add_page_number(slide, 32)


def build_slide33_cdr_results(prs):
    """Slide 33: Phase 2C CDR experiment results & cross-phase summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_blue_accent_bar(slide, 0.0, 0.06)
    add_title_block(slide, "Phase 2C — CDR 实验结果与 Phase 2 总结",
                    "CDR 二分类 AUC + 生存 C-index  |  CDR vs DXSUM 跨终点对比  |  Phase 2 核心发现")

    # Left: CDR Binary results
    add_table(slide, 0.6, 1.4,
              [1.3, 0.9, 0.9, 1.4, 0.7, 1.0, 0.8],
              ["队列", "窗口", "N", "AUC±Std", "Brier", "Sens", "Spec"],
              [["MCI", "3yr", "704", "0.840±0.024", "0.153", "0.847", "0.701"],
               ["MCI", "5yr", "576", "0.886±0.035", "0.144", "0.782", "0.862"],
               ["MCI", "10yr", "444", "0.938±0.021", "0.073", "0.898", "0.892"],
               ["CN", "3yr", "546", "0.705±0.056", "0.105", "0.896", "0.438"],
               ["CN", "5yr", "439", "0.738±0.072", "0.170", "0.592", "0.781"],
               ["CN", "10yr", "281", "0.789±0.038", "0.168", "0.720", "0.761"]])

    # Right: CDR Survival results
    add_table(slide, 5.5, 1.4,
              [1.3, 1.2, 0.9, 0.9, 1.6, 1.5, 1.5],
              ["队列", "终点", "N", "事件", "C-index±Std", "tAUC@3yr", "tAUC@5yr"],
              [["MCI", "Primary", "962", "373", "0.781±0.019", "0.826", "0.865"],
               ["MCI", "Sustained", "962", "223", "0.799±0.029", "0.833", "0.864"],
               ["MCI", "CDRSB≥1", "962", "620", "0.677±0.034", "0.759", "0.793"],
               ["MCI", "Composite", "960", "500", "0.744±0.011", "0.817", "0.848"],
               ["CN", "Primary", "633", "208", "0.680±0.025", "0.649", "0.681"],
               ["CN", "CDRSB≥1", "633", "146", "0.692±0.019", "0.676", "0.683"]])

    # Bottom: CDR vs DXSUM comparison
    add_table(slide, 0.6, 4.45,
              [1.8, 2.2, 2.2, 2.2, 2.2],
              ["指标", "DXSUM MCI", "CDR MCI", "Δ(CDR−DXSUM)", "结论"],
              [["二分类 AUC@3yr", "0.834±0.032", "0.840±0.024", "+0.006", "CDR 略高"],
               ["二分类 AUC@5yr", "0.855±0.021", "0.886±0.035", "+0.031", "CDR 更高"],
               ["生存 C-index", "0.765±0.009", "0.781±0.019", "+0.016", "CDR 略优"],
               ["CN 生存 C-index", "N/A", "0.680±0.025", "—", "CN 可预测"],
               ["跨终点 #1 特征", "PTWORKHS", "FS_ST29SV/FS_ST88SV", "—", "MRI皮层下体积"]])

    # Phase 2 summary conclusion
    add_bottom_conclusion(slide, 6.2, [
        ("Phase 2 核心发现:", ""),
        ("①", " CDR MCI 性能略超 DXSUM  |  "),
        ("②", " CN→MCI 可预测 (C=0.68), Tau PET 是极早期最强信号  |  "),
        ("③", " MRI皮层下体积 是跨终点、跨队列的统一#1特征  |  "),
        ("④", " 生存模型用 1 个模型替代 3 个二分类, 保留全部样本, 自动一致"),
    ])

    add_page_number(slide, 33)


# ── Main ───────────────────────────────────────────────────

def main():
    print(f"Opening: {PPTX_SRC}")
    prs = Presentation(PPTX_SRC)
    orig_count = len(prs.slides)
    print(f"Existing slides: {orig_count}")

    # Build all 7 new slides
    builders = [
        build_slide27_section,
        build_slide28_overview,
        build_slide29_phase2a,
        build_slide30_survival_model,
        build_slide31_survival_features,
        build_slide32_cdr_labels_features,
        build_slide33_cdr_results,
    ]

    for builder in builders:
        builder(prs)
        print(f"  ✓ {builder.__name__}")

    new_count = len(prs.slides)
    print(f"\nTotal slides: {new_count} (added {new_count - orig_count})")

    # Save
    prs.save(PPTX_OUT)
    print(f"Saved: {PPTX_OUT}")


if __name__ == "__main__":
    main()
