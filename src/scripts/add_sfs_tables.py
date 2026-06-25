#!/usr/bin/env python3
"""
Insert a new slide with three complete SFS Top-10 feature tables
between slides 12 and 13 in UKB-DRP_完整实验报告.pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

PPTX = "docs/PPT/UKB-DRP_完整实验报告.pptx"
BLUE = RGBColor(0x00, 0x66, 0xCC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_FILL = RGBColor(0xE8, 0xF0, 0xFE)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
FONT = "Helvetica"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


def set_font(run, size_pt, bold=False, color=GRAY):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def add_blue_bar(slide, top_in=0.0, height_in=0.06):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(top_in), SLIDE_W, Inches(height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def add_title(slide, title, subtitle):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.45))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title; set_font(r, 25, bold=True, color=DARK)

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.5), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = subtitle; set_font(r2, 12, bold=False, color=GRAY)


def add_page_num(slide, num):
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{num}/17"; set_font(r, 7, bold=False, color=GRAY)


def add_sfs_table(slide, left_in, top_in, col_widths_in, title, headers, rows):
    """Compact SFS table with title label above."""
    # Title label
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in - 0.22), Inches(4.0), Inches(0.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"▎{title}"; set_font(r, 8, bold=True, color=BLUE)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths_in)
    row_h = 0.185  # compact row height

    shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left_in), Inches(top_in),
        Inches(total_w), Inches(row_h * n_rows),
    )
    table = shape.table
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    # Header row
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(hdr); set_font(r, 6.5, bold=True, color=WHITE)
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': '0066CC'})
        sf.append(sc); tcPr.append(sf)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            set_font(r, 6.5, bold=(j == 0), color=DARK)
            # Alternating row bg
            if i % 2 == 1:
                tcPr = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F5F7FA'})
                sf.append(sc); tcPr.append(sf)

    return shape


def add_bottom_note(slide, top_in, lines):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_in), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = text; set_font(r, 9, bold=(i == 0), color=BLUE if i == 0 else DARK)
    return txBox


def main():
    prs = Presentation(PPTX)

    # We need to insert after slide 12. python-pptx doesn't support insert,
    # so we rebuild the slide list and reorder.
    # Simpler approach: add at end, then provide instructions to reorder.
    # Actually: add as new slide 13, push old 13 to 14.

    # Build the new slide
    # Use blank layout from existing master
    blank_layout = prs.slide_layouts[6]  # Blank

    # ── Create NEW slide ──
    new_slide = prs.slides.add_slide(blank_layout)
    add_blue_bar(new_slide, 0.0, 0.06)
    add_title(new_slide,
              "Phase 2C — CDR 特征选择: 完整 SFS Top 10 (二分类 5yr + 生存)",
              "MCI 5yr  |  CN 5yr  |  MCI 生存 Primary  —  每张表完整 10 步特征选择")

    # ── Table 1: MCI Primary 5yr SFS ──
    hdr1 = ["#", "特征", "AUC", "增益", "类型"]
    rows1 = [
        ["1", "FS_ST88SV",          "0.771", "+0.771", "MRI 皮层下体积"],
        ["2", "FS_ST31TA",          "0.811", "+0.040", "MRI 皮层厚度"],
        ["3", "APOE4_count",        "0.833", "+0.022", "Bio (遗传)"],
        ["4", "FS_ST40CV",          "0.857", "+0.024", "MRI 皮层下体积"],
        ["5", "FS_ST44TA",          "0.869", "+0.012", "MRI 皮层厚度"],
        ["6", "pT217_AB42_F",       "0.886", "+0.016", "Bio (血浆 pTau217/Aβ42)"],
        ["7", "WMH_LEFT_HIPPO",     "0.888", "+0.002", "WMH"],
        ["8", "FS_ST125SV",         "0.888", "0.000",  "MRI 皮层下体积"],
        ["9", "FS_ST85TS",          "0.889", "+0.001", "MRI 厚度 SD"],
        ["10","FS_ST30SV",          "0.893", "+0.004", "MRI 皮层下体积"],
    ]
    add_sfs_table(new_slide, 0.4, 1.35, [0.22, 2.2, 0.6, 0.55, 1.35], "MCI Primary 5yr 二分类", hdr1, rows1)

    # ── Table 2: CN Primary 5yr SFS ──
    hdr2 = ["#", "特征", "AUC", "增益", "类型"]
    rows2 = [
        ["1", "FS_ST58TA",          "0.606", "+0.606", "MRI 皮层厚度"],
        ["2", "FS_ST90TS",          "0.627", "+0.021", "MRI 厚度 SD"],
        ["3", "FS_ST7SV",           "0.675", "+0.048", "MRI 皮层下体积"],
        ["4", "AB40_F",             "0.701", "+0.026", "Bio (血浆 Aβ40)"],
        ["5", "TAU_META_TEMP_SUVR", "0.711", "+0.010", "Tau PET"],
        ["6", "FS_ST74CV",          "0.722", "+0.010", "MRI 皮层下体积"],
        ["7", "FS_ST76SV",          "0.720", "−0.002", "MRI 皮层下体积"],
        ["8", "NfL_F",              "0.734", "+0.014", "Bio (血浆 NfL)"],
        ["9", "PTEDUCAT",           "0.735", "+0.001", "Bio (教育年限)"],
        ["10","FS_ST119TA",         "0.735", "−0.000", "MRI 皮层厚度"],
    ]
    add_sfs_table(new_slide, 4.7, 1.35, [0.22, 2.4, 0.6, 0.55, 1.35], "CN Primary 5yr 二分类", hdr2, rows2)

    # ── Table 3: MCI Primary 生存 SFS ──
    hdr3 = ["#", "特征", "C-index", "增益", "类型"]
    rows3 = [
        ["1", "FS_ST88SV",          "0.687", "+0.687", "MRI 皮层下体积"],
        ["2", "APOE4_count",        "0.725", "+0.038", "Bio (遗传)"],
        ["3", "FS_ST50TA",          "0.743", "+0.019", "MRI 皮层厚度"],
        ["4", "FS_ST89SV",          "0.759", "+0.016", "MRI 皮层下体积"],
        ["5", "FS_ST24CV",          "0.772", "+0.012", "MRI 皮层下体积"],
        ["6", "GFAP_F",             "0.784", "+0.013", "Bio (血浆 GFAP)"],
        ["7", "FS_ST90TA",          "0.791", "+0.007", "MRI 皮层厚度"],
        ["8", "AMY_SUMMARY_SUVR",   "0.795", "+0.004", "Amyloid PET"],
        ["9", "FS_ST99CV",          "0.800", "+0.005", "MRI 皮层下体积"],
        ["10","FS_ST73TA",          "0.803", "+0.003", "MRI 皮层厚度"],
    ]
    add_sfs_table(new_slide, 9.0, 1.35, [0.22, 2.0, 0.65, 0.55, 1.35], "MCI Primary 生存 (C-index)", hdr3, rows3)

    # ── Bottom note ──
    add_bottom_note(new_slide, 3.85, [
        "▎跨窗口与跨终点特征一致性",
        "  MCI 5yr #1 = FS_ST88SV (皮层下体积) — 与 MCI 3yr #2 (FS_ST29SV) 同为皮层下体积，证实该模态跨窗口稳定",
        "  CN 5yr 起步 AUC 仅 0.606 — CN 队列的长期预测难度极高，第4步才引入第一个 Bio 标志物 (Aβ40)",
        "  MCI 生存 SFS 与 MCI all-time 二分类 SFS 完全相同 — 生存和二分类在同队列上选出完全一致的特征",
        "  MRI 皮层下体积 跨队列 (CN+MCI)、跨终点 (DXSUM+CDR)、跨任务 (二分类+生存) 均为 #1 特征 — 最稳健的 AD 进展影像标志物",
    ])

    add_page_num(new_slide, 16)

    # ── Reorder via XML: move last slide (new) to position 12 (after slide 12) ──
    sldIdLst = prs.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
    if sldIdLst is not None:
        items = list(sldIdLst)
        new_item = items[-1]
        sldIdLst.remove(new_item)
        # Insert after the 12th element (index 12 in 0-based, so before the 13th)
        sldIdLst.insert(12, new_item)
        print(f"  Reordered: new slide moved to position 13")

    # Fix page numbers: each slide with "num/total" pattern
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '/' in run.text and run.text.replace('/', '').replace(str(total-1), '').strip().isdigit():
                            run.text = f"{page_num}/{total}"
                            break

    prs.save(PPTX)
    print(f"✓ Updated {PPTX}")
    print(f"  New slide with 3 SFS Top-10 tables inserted at position 13")
    print(f"  Total slides: {total}")


if __name__ == "__main__":
    main()
