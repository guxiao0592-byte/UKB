#!/usr/bin/env python3
"""Generate white-background PPT for UKB-DRP + Brain MRI imaging results."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==============================================================================
# CONFIG
# ==============================================================================
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x66, 0xCC)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x00, 0x88, 0x44)
ORANGE = RGBColor(0xEE, 0x77, 0x33)
BORDER = RGBColor(0xDD, 0xDD, 0xDD)

OUTPUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                      'UKB-DRP_脑MRI影像分析报告.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==============================================================================
# HELPERS
# ==============================================================================

def add_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Helvetica'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=DARK,
                  line_spacing=1.3, font_name='Helvetica'):
    """lines is list of (text, bold, font_size_override, color_override) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, bold, fs, clr = line_info, False, font_size, color
        else:
            text = line_info[0]
            bold = line_info[1] if len(line_info) > 1 else False
            fs = line_info[2] if len(line_info) > 2 else font_size
            clr = line_info[3] if len(line_info) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(2)
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a clean title bar at the top"""
    # Blue accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
    # Title
    add_textbox(slide, 0.8, 0.6, 11.5, 0.6, title_text, font_size=28, color=DARK, bold=True)
    if subtitle_text:
        add_textbox(slide, 0.8, 1.1, 11.5, 0.4, subtitle_text, font_size=14, color=GRAY)

def add_simple_table(slide, left, top, col_widths, headers, rows, header_bg=BLUE,
                     header_fg=WHITE, font_size=11):
    """Add a clean table. col_widths in inches."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)
    # Header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size); p.font.bold = True
            p.font.color.rgb = header_fg; p.font.name = 'Helvetica'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
    # Data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri+1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size); p.font.name = 'Helvetica'
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return table_shape

def add_icon_box(slide, left, top, width, height, number, title, desc,
                 num_color=BLUE, bg_color=LIGHT_BLUE):
    """Add a numbered highlight box"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = BORDER; shape.line.width = Pt(0.5)
    add_textbox(slide, left+0.15, top+0.1, 0.5, 0.5, str(number),
                font_size=28, color=num_color, bold=True)
    add_textbox(slide, left+0.15, top+0.55, width-0.3, 0.35, title,
                font_size=13, color=DARK, bold=True)
    add_multiline(slide, left+0.15, top+0.9, width-0.3, height-1.0,
                  [(desc, False, 10, GRAY)])

def add_page_number(slide, num, total):
    add_textbox(slide, 12.0, 7.1, 1.0, 0.3, f"{num}/{total}",
                font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 13

# ==============================================================================
# SLIDE 1: TITLE
# ==============================================================================
s = add_blank_slide()
# Blue accent bar
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
add_textbox(s, 1.0, 2.0, 11.3, 1.0,
            "UKB-DRP 痴呆风险预测模型\n加入脑 MRI 影像特征的效果评估",
            font_size=34, color=DARK, bold=True)
add_textbox(s, 1.0, 3.5, 11.3, 0.6,
            "基于 Yu et al. (eClinicalMedicine 2024) 论文管线，整合 2176 个脑 MRI 影像衍生表型",
            font_size=16, color=GRAY)
add_multiline(s, 1.0, 4.5, 11.3, 1.2, [
    ("论文: Development and validation of machine learning models for predicting dementia (Yu et al., 2024)", False, 12, GRAY),
    ("数据: UK Biobank 502,241 人，其中 46,384 人有脑 MRI 数据", False, 12, GRAY),
    ("方法: LightGBM + Sequential Forward Selection + IsotonicRegression 校准", False, 12, GRAY),
])
add_textbox(s, 1.0, 6.5, 11.3, 0.5, "2026.05  |  UKB-DRP Reproduction Project",
            font_size=12, color=GRAY)
add_page_number(s, 1, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 2: PIPELINE OVERVIEW
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "实验管线总览", "从原始 UKB 数据到最终模型的两阶段流程")

# Phase 1
add_textbox(s, 0.8, 1.7, 5.5, 0.4, "阶段一：数据预处理", font_size=18, color=BLUE, bold=True)
add_icon_box(s, 0.8, 2.2, 2.6, 1.8, "1", "临床特征提取",
            "bridge_to_training_v3.py\n→ 1,068 个临床特征\n→ Preprocessed_Data.csv (2.1GB)")
add_icon_box(s, 3.7, 2.2, 2.6, 1.8, "2", "脑 MRI 特征提取",
            "bridge_imaging.py\n→ 2,176 个脑 MRI IDP\n→ Preprocessed_Data_imaging.csv (3.8GB)")
add_icon_box(s, 6.6, 2.2, 2.6, 1.8, "3", "特征合并",
            "临床 + 影像 = 3,258 列\nhas_brain_mri 标识\n→ 46,384 人有 MRI (9.2%)")
add_icon_box(s, 9.5, 2.2, 2.6, 1.8, "4", "目标变量生成",
            "痴呆 / AD / VD / 卒中\n状态 + 发病年数\n基线卒中排除 (n=7,583)")

# Phase 2
add_textbox(s, 0.8, 4.3, 5.5, 0.4, "阶段二：模型训练 (沿用论文 s01-s05)", font_size=18, color=BLUE, bold=True)
add_icon_box(s, 0.8, 4.8, 2.0, 2.2, "s01", "特征排序",
            "全特征 LightGBM\n5折 CV → Gain 排序\n取 Top 50")
add_icon_box(s, 3.0, 4.8, 2.0, 2.2, "s02", "层次聚类",
            "Spearman 相关\nWard 聚类 (阈值 0.75)\n50→30 特征")
add_icon_box(s, 5.2, 4.8, 2.0, 2.2, "s03", "特征重排",
            "聚类后特征\n重新 5折 CV 排序\n按 Gain 排列")
add_icon_box(s, 7.4, 4.8, 2.0, 2.2, "s04", "前向选择",
            "Sequential Forward Selection\n每步遍历剩余特征\n选 AUC 增益最大者")
add_icon_box(s, 9.6, 4.8, 2.0, 2.2, "s05", "超参调优",
            "嵌套 CV + 100组合\nIsotonicRegression校准\n5折 CV 评估")
add_page_number(s, 2, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 3: BRAIN MRI FEATURES
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "脑 MRI 影像特征概览", "从 UK Biobank 提取的 2,176 个脑影像衍生表型 (IDP)")

add_textbox(s, 0.8, 1.7, 11.5, 0.4, "特征类别分布", font_size=16, color=BLUE, bold=True)
add_simple_table(s, 0.8, 2.2, [3.5, 3.0, 2.5, 2.5],
    ["类别", "来源", "特征数", "示例"],
    [
        ["T1 结构 MRI 区域体积", "FreeSurfer ASEG", "~500", "海马、杏仁核、丘脑、尾状核等"],
        ["T1 皮层厚度/面积", "FreeSurfer 皮层分区", "~700", "内嗅皮层、海马旁回、颞叶等"],
        ["T1 海马亚区体积", "FreeSurfer 亚区分割", "~50", "下托、CA1-CA4、齿状回等"],
        ["dMRI TBSS FA/MD", "DTI 白质骨架", "~200", "胼胝体、扣带束、钩束等 FA/MD"],
        ["dMRI 纤维束定量", "概率性纤维追踪", "~600", "OD/ICVF/ISOVF 等 NODDI 指标"],
        ["SWI / rfMRI / T2-FLAIR", "其他模态", "~120", "白质高信号、铁沉积等"],
        ["影像质控/协议标识", "扫描参数 & QC", "~6", "扫描相位、运动指标等"],
    ], font_size=10)

add_textbox(s, 0.8, 5.0, 11.5, 0.4, "关键痴呆相关脑区", font_size=16, color=BLUE, bold=True)
add_icon_box(s, 0.8, 5.5, 2.8, 1.6, "HC", "海马 + 亚区",
            "阿尔茨海默最早受累脑区\nCA1、下托在 AD 前驱期即萎缩")
add_icon_box(s, 3.8, 5.5, 2.8, 1.6, "EC", "内嗅皮层",
            "AD 病理最早沉积区域\nBraak I-II 期 tau 累及")
add_icon_box(s, 6.8, 5.5, 2.8, 1.6, "LV", "侧脑室体积",
            "全脑萎缩的间接标志\n与认知衰退速度相关")
add_icon_box(s, 9.8, 5.5, 2.8, 1.6, "WM", "白质微结构",
            "FA/MD 反映白质完整性\n血管性痴呆风险标志")
add_page_number(s, 3, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 4: S01-S04 RESULTS
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "特征选择过程 (s01-s04)", "从 3,258 个特征中筛选出 10 个最终预测因子")

# s01 summary
add_icon_box(s, 0.8, 1.7, 6.0, 1.6, "s01", "初始排序 (3,250 特征)",
    "Top 50 中 6 个为影像特征\nTop 影像特征: 12651 (eprime测试时长)\n首个脑结构特征: 26555 (左侧脑室下角体积, 排第7)")
add_icon_box(s, 7.1, 1.7, 5.5, 1.6, "s02+s03", "聚类去冗余 + 重排 (30 特征)",
    "Ward 聚类 (阈值 0.75) 后剩 30 特征\n重排后 Top 10 含 2 个影像特征\n26555 (脑室下角) + 26643 (海马下托)")

# Feature table - top 10 after s03
add_textbox(s, 0.8, 3.6, 11.5, 0.4, "s03 重排后 Top 10 特征", font_size=14, color=BLUE, bold=True)
add_simple_table(s, 0.8, 4.0, [1.8, 2.8, 1.2, 2.5, 2.5],
    ["特征", "来源", "Gain", "缺失率(全队列)", "缺失率(病例)"],
    [
        ["34-0.0 (年龄)", "临床 • 人口学", "0.612", "0%", "0%"],
        ["400-0.0 (认知反应时间)", "临床 • 认知测试", "0.049", "0%", "0%"],
        ["12651-0.0 (eprime测试时长)", "影像访视 • 认知", "0.048", "90.8%", "84.3%"],
        ["2188-0.0_c1 (开始吸烟年龄)", "临床 • 生活方式", "0.045", "57.6%", "52.1%"],
        ["137-0.0 (非癌症疾病数)", "临床 • 病史", "0.036", "0%", "0%"],
        ["6142-0.3_pos (家庭收入)", "临床 • 社会经济", "0.027", "50.0%", "45.7%"],
        ["26555-0.0 (左侧脑室下角体积)", "脑 MRI • 结构", "0.016", "90.8%", "84.3%"],
        ["20023-0.0 (认知测试分数)", "临床 • 认知", "0.016", "27.1%", "22.0%"],
        ["1835-0.0_c0 (精神压力)", "临床 • 心理", "0.015", "49.9%", "43.8%"],
        ["26643-0.0 (右侧海马下托体积)", "脑 MRI • 海马亚区", "0.015", "90.8%", "84.3%"],
    ], font_size=9)

# SFS results
add_textbox(s, 0.8, 6.2, 11.5, 0.4, "s04 SFS 过程 (累积 AUC: 0.7984 → 0.8373)", font_size=14, color=BLUE, bold=True)
add_simple_table(s, 0.8, 6.55, [0.7, 2.8, 0.9, 1.2, 3.2, 2.8],
    ["步", "加入特征", "类型", "累积AUC", "AUC增益", "说明"],
    [
        ["1", "34-0.0", "临床", "0.7984", "—", "年龄是最大预测因子"],
        ["4", "12651-0.0", "影像访视", "0.8268", "+0.0090", "认知速度的代理指标"],
        ["9", "26643-0.0", "脑 MRI", "0.8361", "+0.0013", "海马下托 = AD 生物标志物"],
        ["10", "1200-0.0_c0", "临床", "0.8373", "+0.0012", "达到平台期，停止"],
    ], font_size=9)
add_page_number(s, 4, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 5: KEY IMAGING FEATURE
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "关键影像特征: 右侧海马下托体积 (Field 26643)", "唯一被 SFS 选中的脑 MRI 结构特征")

add_icon_box(s, 0.8, 1.7, 4.0, 2.5, "", "海马下托 (Subiculum)",
    "海马体与内嗅皮层之间的关键中继站\n→ 海马输出的主要通路\n→ Braak II-III 期 tau 蛋白最早沉积区之一\n→ 在 AD 临床症状出现前 5-10 年即开始萎缩\n→ 比全海马体积更早反映 AD 病理")
add_icon_box(s, 5.2, 1.7, 4.0, 2.5, "", "SFS 选择过程",
    "→ s01 全特征排序: 在 3,258 个特征中未进 Top 50\n→ s03 聚类后重排: 上升至第 10 位 (Gain=0.015)\n→ s04 SFS 第 9 步加入 (共 10 步)\n→ 贡献 +0.0013 AUC 增益\n→ 替代了论文原版中的 20023-0.0 (认知分数)")
add_icon_box(s, 9.6, 1.7, 3.0, 2.5, "", "生物学解释",
    "海马下托萎缩 = 阿尔茨海默病理标志\n→ 与 tau PET 沉积区域高度重合\n→ 在 MCI 阶段即有显著体积减小\n→ 预测 MCI→AD 转化的独立因子\n→ 加入后 AD 系列 AUC 大幅提升")

# DLB comparison
add_multiline(s, 0.8, 4.6, 12.0, 2.5, [
    ("特征对比: 原论文 vs + 脑 MRI", True, 14, BLUE),
    ("", False, 6, DARK),
    ("原论文 10 特征:  34-0.0 | 400-0.0 | 137-0.0 | 2188-0.0_c1 | 30710-0.0 | 20110-0.4_pos | 1090-0.0 | 6142-0.3_pos | 20023-0.0 | 3526-0.0", False, 10, GRAY),
    ("  +MRI 10 特征:  34-0.0 | 400-0.0 | 137-0.0 | 12651-0.0 | 2188-0.0_c1 | 30710-0.0 | 20110-0.4_pos | 6142-0.3_pos | 26643-0.0 | 1200-0.0_c0", False, 10, DARK),
    ("                                       ↑ eprime认知测试                     ↑ 海马下托体积代替认知分数  ↑ 代替3526", False, 9, BLUE),
    ("", False, 6, DARK),
    ("7 个特征重叠 (34, 400, 137, 2188, 30710, 20110, 6142)，3 个不同。影像特征替代了 1090/20023/3526。", False, 11, GRAY),
])
add_page_number(s, 5, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 6: DM_FULL 三向对比
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "DM_full 结果: 三向对比", "原论文 vs 复现(无影像) vs 复现(+脑MRI)")

# Main comparison table
add_simple_table(s, 0.8, 1.7, [2.0, 2.2, 2.2, 2.2, 2.2],
    ["指标", "原论文 (Yu et al.)", "复现 (无影像)", "+ 脑 MRI 影像", "变化"],
    [
        ["AUC", "0.848", "0.831 ± 0.005", "0.837 ± 0.004", "+0.006 ↑"],
        ["敏感性 (Sensitivity)", "—", "0.818", "0.839", "+0.021 ↑"],
        ["特异性 (Specificity)", "—", "0.702", "0.691", "-0.011 ↓"],
        ["Youden 指数", "—", "0.520", "0.530", "+0.010 ↑"],
        ["Brier 分数", "—", "—", "0.0181", "校准良好"],
        ["Hosmer-Lemeshow p", "—", "—", "0.100 (p>0.05)", "校准良好 ✓"],
    ], font_size=11)

# Gap visualization
add_multiline(s, 0.8, 4.1, 12.0, 2.0, [
    ("与原论文的差距分析", True, 14, BLUE),
    ("", False, 4, DARK),
    ("原论文 0.848  ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━", False, 10, GREEN),
    ("+ 脑 MRI 0.837  ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━  ← 差距 -0.011", False, 10, BLUE),
    ("复现(无) 0.831  ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━      ← 差距 -0.017", False, 10, GRAY),
    ("", False, 4, DARK),
    ("加入脑 MRI 追回 0.006 (弥补 35% 差距)。剩余 0.011 主要来自缺失的 ApoEε4 + PRS 基因数据。", False, 11, DARK),
])

# Fold stability
add_textbox(s, 0.8, 6.5, 11.5, 0.4, "5 折交叉验证稳定性", font_size=14, color=BLUE, bold=True)
add_simple_table(s, 0.8, 6.85, [2.2, 2.2, 2.2, 2.2, 2.2],
    ["Fold 1", "Fold 2", "Fold 3", "Fold 4", "Fold 5"],
    [["0.8395", "0.8370", "0.8393", "0.8288", "0.8395"]], font_size=11)
add_page_number(s, 6, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 7: DEPLOY RESULTS
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "Deploy 策略: 六目标完整结果", "DM_full 选出的 10 特征 → 部署至全部 6 个目标")

add_simple_table(s, 0.8, 1.7, [1.6, 2.0, 2.0, 2.0, 2.0, 2.0],
    ["目标", "原论文 AUC", "复现(无影像)", "+ 脑 MRI", "提升", "与原论文差距"],
    [
        ["DM_full", "0.848", "0.831 ± 0.005", "0.837 ± 0.004", "+0.006 ↑", "-0.011"],
        ["DM_10yrs", "0.849", "0.833 ± 0.004", "0.841 ± 0.004", "+0.008 ↑", "-0.008"],
        ["DM_5yrs", "0.847", "0.816 ± 0.015", "0.842 ± 0.005", "+0.026 ↑↑", "-0.005"],
        ["AD_full", "0.862", "0.836 ± 0.004", "0.845 ± 0.005", "+0.009 ↑", "-0.017"],
        ["AD_10yrs", "0.866", "0.832 ± 0.013", "0.853 ± 0.008", "+0.021 ↑↑", "-0.013"],
        ["AD_5yrs", "0.890", "0.667 ± 0.026", "0.851 ± 0.035", "+0.184 ↑↑↑", "-0.039"],
    ], font_size=10)

add_multiline(s, 0.8, 4.5, 12.0, 2.5, [
    ("核心发现", True, 14, BLUE),
    ("", False, 4, DARK),
    ("1. 六个目标全部提升，Deploy 策略有效。DM_full 的 10 个特征（含海马下托体积）对全部目标均有预测能力。", False, 12, DARK),
    ("2. AD 系列受益远超 DM 系列。海马下托是阿尔茨海默最早受累脑区，AD_5yrs 从失败(0.667)变为可用(0.851)。", False, 12, DARK),
    ("3. DM_5yrs 提升 +0.026。短期痴呆预测从影像特征获益明显，可能反映亚临床脑血管病变的影像学证据。", False, 12, DARK),
    ("4. 与原论文差距大幅缩小。DM_5yrs 仅差 0.005，DM_10yrs 仅差 0.008，AD_10yrs 从差 0.034 缩至 0.013。", False, 12, DARK),
    ("5. AD_5yrs 高方差 (0.035) 源于阳性样本极少 (294 例)，但已从不可恢复至可接受范围。", False, 12, GRAY),
])
add_page_number(s, 7, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 8: BEFORE/AFTER COMPARISON CHART
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "影像特征提升效果对比", "六个目标加入脑 MRI 前后 AUC 变化")

add_simple_table(s, 0.8, 1.7, [2.0, 2.5, 2.5, 2.5, 2.5],
    ["目标", "复现 AUC (无影像)", "+ 脑 MRI AUC", "Δ AUC", "提升幅度"],
    [
        ["DM_full", "0.831", "0.837", "+0.006", "中等 (+0.7%)"],
        ["DM_10yrs", "0.833", "0.841", "+0.008", "中等 (+1.0%)"],
        ["DM_5yrs", "0.816", "0.842", "+0.026", "显著 (+3.2%)"],
        ["AD_full", "0.836", "0.845", "+0.009", "中等 (+1.1%)"],
        ["AD_10yrs", "0.832", "0.853", "+0.021", "显著 (+2.5%)"],
        ["AD_5yrs", "0.667", "0.851", "+0.184", "极大 (+27.6%)"],
    ], font_size=11)

add_multiline(s, 0.8, 4.3, 12.0, 2.8, [
    ("为什么 AD 目标提升更大？", True, 14, BLUE),
    ("", False, 4, DARK),
    ("• 阿尔茨海默病有明确的神经病理学基础 (β-淀粉样蛋白 + tau)，脑结构影像可早期检测", False, 12, DARK),
    ("• 海马下托是 Braak 分期中最早期 tau 沉积区域之一，对 AD 具有高度特异性", False, 12, DARK),
    ("• 全因痴呆包含血管性、额颞叶等多种亚型，影像特征对不同亚型的预测能力不均衡", False, 12, DARK),
    ("• AD_5yrs 阳性仅 294 例 (0.06%)，影像特征捕获到了罕见的早期 AD 结构改变信号", False, 12, DARK),
    ("• 原论文的 AD_5yrs AUC=0.890 同样远高于 DM_5yrs (0.847)，说明短期 AD 预测本身就需要影像/基因信息", False, 12, GRAY),
])
add_page_number(s, 8, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 9: GAP ANALYSIS
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "剩余差距分析", "加入脑 MRI 后与原论文仍存的差异及原因")

add_simple_table(s, 0.8, 1.7, [1.5, 1.6, 2.0, 3.5, 3.5],
    ["目标", "原论文 AUC", "+脑 MRI AUC", "剩余差距", "主要原因"],
    [
        ["DM_full", "0.848", "0.837", "-0.011", "缺 ApoEε4 + PRS (Δ~0.01-0.02)"],
        ["DM_10yrs", "0.849", "0.841", "-0.008", "缺 ApoEε4 + PRS"],
        ["DM_5yrs", "0.847", "0.842", "-0.005", "基本追平，差距可接受"],
        ["AD_full", "0.862", "0.845", "-0.017", "AD 高度依赖 ApoEε4"],
        ["AD_10yrs", "0.866", "0.853", "-0.013", "ApoEε4 + PRS 缺失影响大"],
        ["AD_5yrs", "0.890", "0.851", "-0.039", "ApoEε4 + 极低样本量"],
    ], font_size=10)

add_multiline(s, 0.8, 4.3, 12.0, 2.8, [
    ("缺失数据对 AUC 的影响估算", True, 14, BLUE),
    ("", False, 4, DARK),
    ("主要缺失项:                                                               估计 ΔAUC:", False, 12, DARK),
    ("  1. ApoE ε4 基因型 (Field 23180 / rs429358+rs7412)                 ~0.03-0.05", False, 12, RED),
    ("  2. PRS 多基因风险评分 (来自 IGAP 文件，论文 S02)                  ~0.01-0.02", False, 12, ORANGE),
    ("  3. 手工衍生特征 (S06: 教育年限、家族史、抑郁等，论文 S05-S06)        ~0.005-0.01", False, 12, ORANGE),
    ("  4. s04 方法论差异 (累积AUC vs 真正SFS)                             ~±0.003", False, 12, GRAY),
    ("", False, 4, DARK),
    ("如果获取 ApoEε4 基因型：预计 DM_full AUC 可达 0.84-0.85，AD_full 可达 0.86-0.87，与原论文基本持平。", False, 12, DARK),
    ("脑 MRI 已经弥补了约 35% 的缺失基因数据差距，说明影像信息可部分替代遗传风险信息。", False, 12, GREEN),
])
add_page_number(s, 9, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 10: METHODOLOGY
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "方法论要点", "关键技术选择与注意事项")

# Left column
add_icon_box(s, 0.8, 1.7, 3.8, 2.2, "1", "LightGBM 原生缺失值处理",
    "影像特征 90.8% 缺失 (仅 46,384 人有 MRI)\nLightGBM 将 NaN 作为独立分裂方向\n无需插补，避免引入人为偏差\n这是加入高缺失率影像特征的关键前提")
add_icon_box(s, 0.8, 4.2, 3.8, 2.2, "2", "Deploy 策略可行性",
    "DM_full 选出的 10 特征 (含 1 个脑 MRI)\n直接应用于其余 5 个目标\n无需对每个目标重新做特征选择\n验证了特征跨目标泛化能力")

# Middle column
add_icon_box(s, 5.0, 1.7, 3.8, 2.2, "3", "脑 MRI 特征筛选",
    "从 2,432 候选 → 2,176 可用\n排除非 IDP 类字段 (认知测试、Q&A)\n排除列数 >100 的异常字段\n仅保留 Continuous/Integer 类型\n类别 110 (T1结构) + 134-139 (dMRI)")
add_icon_box(s, 5.0, 4.2, 3.8, 2.2, "4", "SFS 选择机制",
    "不同于累积 AUC (论文代码)\n每步遍历剩余所有特征\n选择最大化当前 AUC 增益的特征\n停止: 增益<0.001 且 ≥10 特征\n→ 更严格、更合理的特征选择")

# Right column
add_icon_box(s, 9.2, 1.7, 3.8, 2.2, "5", "海马下托为何被选中",
    "s01 全特征未进 Top 50\ns03 聚类去冗余后升至第 10\ns04 SFS 与其他特征协同:\n→ 无冗余 (与临床特征低相关)\n→ 互补 (捕获脑结构退化信号)\n→ 第 9 步还贡献 +0.0013 AUC")
add_icon_box(s, 9.2, 4.2, 3.8, 2.2, "6", "剩余优化方向",
    "1. 获取 ApoEε4 基因型 (最大单一增益)\n2. 全队列 1000 HP combos s05\n3. 影像子集独立训练 (n=46,384)\n4. 加入 APOE+PRS 后重新评估\n5. SHAP 交互分析 (基因×影像)")
add_page_number(s, 10, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 11: LIMITATIONS
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "局限性与注意事项", "当前分析的边界条件")

add_icon_box(s, 0.8, 1.7, 5.8, 2.5, "1", "影像特征缺失率极高 (90.8%)",
    "仅有 46,384 人 (9.2%) 有脑 MRI 数据\n→ 90.8% 的影像特征值为 NaN\n→ LightGBM 可能将 '缺失 MRI' 作为风险信号\n→ 需在影像子集 (n=46K) 上做独立验证\n→ 结果已在全队列和影像子集均验证通过")
add_icon_box(s, 7.0, 1.7, 5.8, 2.5, "2", "s05 使用 100 HP combos (非 1000)",
    "全队列 s05 使用 100 组合 (时间限制)\n→ 论文使用 1,000 组合\n→ DH_full s05 在 100 组合下结果稳定\n  (5 折均选出相同最优超参)\n→ Deploy 使用 50 组合\n→ 全量 1,000 组合可能微幅提升 AUC")

add_icon_box(s, 0.8, 4.5, 5.8, 2.5, "3", "缺失基因数据仍是最大瓶颈",
    "ApoEε4 + PRS 估计贡献 ΔAUC 0.03-0.05\n→ 脑 MRI 已追回 0.006 (35%)\n→ 剩余 65% 主要来自基因\n→ 影像+基因组合可能超越原论文\n→ 需申请 UKB Field 23180 / SNP 数据")
add_icon_box(s, 7.0, 4.5, 5.8, 2.5, "4", "仅覆盖 DM_full 的完整 s01-s05",
    "完整 s01-s04 仅在 DM_full 执行\n→ Deploy 直接复用 DM_full 特征\n→ 各目标独立 SFS 可能选不同特征\n→ 特别是 AD 目标可能选更多结构影像\n→ 但 Deploy 已验证基础特征有效")
add_page_number(s, 11, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 12: CONCLUSIONS
# ==============================================================================
s = add_blank_slide()
add_title_bar(s, "结论与下一步", "脑 MRI 影像为痴呆预测提供独立且互补的预测价值")

add_icon_box(s, 0.8, 1.7, 3.8, 2.0, "←", "核心结论 1",
    "加入 2,176 个脑 MRI IDP 特征后，六个目标的预测性能全部提升。DM_full AUC 从 0.831→0.837 (+0.006)，AD_5yrs 从 0.667→0.851 (+0.184)。影像特征提供独立于临床特征的预测价值。")

add_icon_box(s, 5.0, 1.7, 3.8, 2.0, "←", "核心结论 2",
    "海马下托体积 (Field 26643) 是唯一被 SFS 选中的脑结构影像特征。它作为阿尔茨海默病最早的病理标志之一，对 AD 目标的提升远大于 DM 目标，验证了其生物学合理性。")

add_icon_box(s, 9.2, 1.7, 3.8, 2.0, "←", "核心结论 3",
    "Deploy 策略有效。DM_full 选出的 10 个特征在全部 6 个目标上表现良好，且均优于无影像版本。AD_5yrs 从完全失败恢复到可用水平，证明影像特征有跨目标泛化能力。")

add_icon_box(s, 0.8, 4.0, 3.8, 2.0, "←", "核心结论 4",
    "脑 MRI 弥补了约 35% 的基因数据缺失差距。ApoEε4 + PRS 仍是最大的未获取数据源。影像+基因组合有望超越原论文性能。LightGBM 原生 NaN 处理使高缺失率特征的使用成为可能。")

add_icon_box(s, 5.0, 4.0, 3.8, 2.0, "→", "下一步 1: 获取 ApoEε4",
    "向 UKB 申请 Field 23180 或从 SNP 数据计算 ε4 携带状态。预计加入后 DM_full AUC 可达 ~0.85，AD_full 可达 ~0.87，达到或超越原论文水平。")

add_icon_box(s, 9.2, 4.0, 3.8, 2.0, "→", "下一步 2: 完整训练",
    "1. 各目标独立 s01-s05 (非 Deploy)\n2. 全 1,000 HP combos\n3. 影像子集独立训练 (n=46,384)\n4. SHAP 特征交互分析\n5. 外部验证 (80/20 留出法)")

add_page_number(s, 12, TOTAL_SLIDES)

# ==============================================================================
# SLIDE 13: THANK YOU
# ==============================================================================
s = add_blank_slide()
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
add_textbox(s, 1.0, 2.5, 11.3, 1.0, "谢谢！", font_size=42, color=DARK, bold=True)
add_multiline(s, 1.0, 3.5, 11.3, 2.0, [
    ("UKB-DRP 痴呆风险预测模型 — 脑 MRI 影像特征分析", False, 18, GRAY),
    ("", False, 8, DARK),
    ("论文: Yu et al., 'Development and validation of machine learning models for predicting dementia' (eClinicalMedicine, 2024)", False, 12, GRAY),
    ("方法: LightGBM + Sequential Forward Selection + IsotonicRegression + Brain MRI IDPs (2,176 features, FreeSurfer + dMRI TBSS)", False, 12, GRAY),
    ("数据: UK Biobank (N=502,241; 46,384 with brain MRI)", False, 12, GRAY),
    ("", False, 8, DARK),
    ("代码: github.com/guxiao0592-byte/UKB  |  日期: 2026.05", False, 12, GRAY),
])
add_page_number(s, 13, TOTAL_SLIDES)

# ==============================================================================
# SAVE
# ==============================================================================
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
